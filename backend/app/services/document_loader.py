import os
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document as DocxDocument
import logging
from typing import Optional, Tuple, List
from pathlib import Path

from ..core.config import settings

logger = logging.getLogger(__name__)


class DocumentLoader:
    """Класс для загрузки и извлечения текста из различных типов документов"""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._extract_pdf_text,
            '.pptx': self._extract_pptx_text,
            '.docx': self._extract_docx_text,
            '.txt': self._extract_txt_text
        }
    
    def extract_text(self, file_path: str) -> Tuple[str, Optional[int]]:
        """
        Извлечение текста из документа
        
        Returns:
            Tuple[str, Optional[int]]: (extracted_text, page_count)
        """
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        try:
            return self.supported_formats[file_extension](file_path)
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            raise
    
    def _extract_pdf_text(self, file_path: str) -> Tuple[str, int]:
        """Извлечение текста из PDF файла"""
        text_content = []
        
        with fitz.open(file_path) as pdf_doc:
            page_count = len(pdf_doc)
            
            for page_num in range(page_count):
                page = pdf_doc[page_num]
                page_text = page.get_text()
                
                if page_text.strip():  # добавляем только непустые страницы
                    text_content.append(f"--- Страница {page_num + 1} ---\n{page_text}")
        
        full_text = "\n\n".join(text_content)
        
        logger.info(f"Extracted text from PDF: {len(full_text)} characters, {page_count} pages")
        return full_text, page_count
    
    def _extract_pptx_text(self, file_path: str) -> Tuple[str, int]:
        """Извлечение текста из PowerPoint файла"""
        text_content = []
        
        presentation = Presentation(file_path)
        slide_count = len(presentation.slides)
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []
            
            # Извлекаем текст из всех текстовых фреймов на слайде
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                slide_content = f"--- Слайд {slide_num} ---\n" + "\n".join(slide_text)
                text_content.append(slide_content)
        
        full_text = "\n\n".join(text_content)
        
        logger.info(f"Extracted text from PPTX: {len(full_text)} characters, {slide_count} slides")
        return full_text, slide_count
    
    def _extract_docx_text(self, file_path: str) -> Tuple[str, int]:
        """Извлечение текста из Word документа"""
        doc = DocxDocument(file_path)
        
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text.strip())
        
        # Извлекаем текст из таблиц
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(" | ".join(row_text))
        
        full_text = "\n\n".join(text_content)
        page_count = len(doc.sections) if doc.sections else 1
        
        logger.info(f"Extracted text from DOCX: {len(full_text)} characters, {page_count} sections")
        return full_text, page_count
    
    def _extract_txt_text(self, file_path: str) -> Tuple[str, int]:
        """Извлечение текста из текстового файла"""
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        # Подсчитываем примерное количество "страниц" (по 3000 символов)
        page_count = max(1, len(content) // 3000)
        
        logger.info(f"Extracted text from TXT: {len(content)} characters, ~{page_count} pages")
        return content, page_count
    
    def extract_pdf_images(self, file_path: str) -> List[Tuple[bytes, int]]:
        """
        Извлечение изображений из PDF для OCR
        
        Returns:
            List[Tuple[bytes, int]]: список (image_data, page_number)
        """
        images = []
        
        try:
            with fitz.open(file_path) as pdf_doc:
                for page_num in range(len(pdf_doc)):
                    page = pdf_doc[page_num]
                    
                    # Получаем изображение страницы
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # увеличиваем разрешение
                    img_data = pix.tobytes("png")
                    images.append((img_data, page_num + 1))
        
        except Exception as e:
            logger.error(f"Error extracting images from PDF {file_path}: {e}")
            raise
        
        return images
    
    def chunk_text(self, text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
        """
        Разбиение текста на чанки для векторизации
        
        Args:
            text: исходный текст
            chunk_size: размер чанка (по умолчанию из настроек)
            overlap: перекрытие между чанками
        
        Returns:
            List[str]: список чанков текста
        """
        chunk_size = chunk_size or settings.MAX_CHUNK_SIZE
        overlap = overlap or settings.CHUNK_OVERLAP
        
        # Разбиваем текст на предложения
        sentences = text.split('.')
        
        chunks = []
        current_chunk = []
        current_length = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            sentence_length = len(sentence)
            
            # Если добавление предложения превысит лимит чанка
            if current_length + sentence_length > chunk_size and current_chunk:
                # Сохраняем текущий чанк
                chunk_text = '. '.join(current_chunk) + '.'
                chunks.append(chunk_text)
                
                # Начинаем новый чанк с перекрытием
                if overlap > 0 and len(current_chunk) > 1:
                    # Берем последние предложения для перекрытия
                    overlap_sentences = current_chunk[-overlap:]
                    current_chunk = overlap_sentences + [sentence]
                    current_length = sum(len(s) for s in current_chunk)
                else:
                    current_chunk = [sentence]
                    current_length = sentence_length
            else:
                current_chunk.append(sentence)
                current_length += sentence_length
        
        # Добавляем последний чанк
        if current_chunk:
            chunk_text = '. '.join(current_chunk) + '.'
            chunks.append(chunk_text)
        
        logger.info(f"Text split into {len(chunks)} chunks")
        return chunks
    
    def get_file_info(self, file_path: str) -> dict:
        """Получение информации о файле"""
        file_stat = os.stat(file_path)
        
        return {
            "size": file_stat.st_size,
            "extension": Path(file_path).suffix.lower(),
            "name": Path(file_path).name,
            "stem": Path(file_path).stem,
            "created": file_stat.st_ctime,
            "modified": file_stat.st_mtime
        }
    
    def validate_file(self, file_path: str) -> bool:
        """Валидация загружаемого файла"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_info = self.get_file_info(file_path)
        
        # Проверяем размер файла
        if file_info["size"] > settings.MAX_FILE_SIZE:
            raise ValueError(f"File too large: {file_info['size']} bytes")
        
        # Проверяем расширение
        if file_info["extension"] not in settings.ALLOWED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {file_info['extension']}")
        
        return True


# Создаем глобальный экземпляр загрузчика
document_loader = DocumentLoader()